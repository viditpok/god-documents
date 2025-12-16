from __future__ import annotations

import sys
from pathlib import Path

import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT / "src") not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT / "src"))

from macrotones.api import loader  # noqa: E402
from macrotones.data.macro_loader import get_macro_features  # noqa: E402
from macrotones.diagnostics.attribution import factor_attribution  # noqa: E402
from macrotones.fusion.llm_allocator import _deterministic_policy  # noqa: E402

SENS_PRESET = {"CPI": 1.0, "UNRATE": -0.5, "T10Y2Y": 0.2}
MACRO_KEY_MAP = {"CPI": "cpi_yoy", "UNRATE": "unemp", "T10Y2Y": "y10y_2y_spread"}


def _save_plot(fig, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.write_image(path, scale=2, width=900, height=520)
    return path


def _add_footer(slide, footer_text: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
    text_frame = textbox.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = footer_text
    font = run.font
    font.size = Pt(10)
    font.bold = False


def _ic_chart(processed: Path) -> Path | None:
    ic_path = processed / "ic_summary.csv"
    if not ic_path.exists():
        return None
    df = pd.read_csv(ic_path)
    if "factor" not in df.columns:
        return None
    df["significant"] = df.get("significant", df.get("pearson_pvalue", 1.0) < 0.05)
    fig = px.bar(
        df,
        x="factor",
        y="pearson_ic",
        color="significant",
        color_discrete_map={True: "#22c55e", False: "#94a3b8"},
        title="Information Coefficients",
    )
    fig.update_layout(xaxis_title="Factor", yaxis_title="Pearson IC")
    return _save_plot(fig, processed / "ic_summary.png")


def _attribution_chart(processed: Path) -> Path | None:
    bt_path = processed / "backtest.parquet"
    if not bt_path.exists():
        return None
    bt = pd.read_parquet(bt_path)
    returns = bt.filter(like="ret_")
    weights = bt.filter(like="w_")
    if returns.empty or weights.empty:
        return None
    returns.columns = [col.replace("ret_", "") for col in returns.columns]
    weights.columns = [col.replace("w_", "") for col in weights.columns]
    overlap = sorted(set(returns.columns) & set(weights.columns))
    if not overlap:
        return None
    attr = factor_attribution(returns[overlap], weights[overlap]).reset_index()
    attr = attr.rename(columns={"index": "Year"})
    fig = px.bar(
        attr,
        x="Year",
        y="sum",
        color="sum",
        color_continuous_scale="RdBu",
        title="Yearly Attribution",
    )
    fig.update_layout(yaxis_title="Contribution")
    return _save_plot(fig, processed / "attribution.png")


def _quintile_chart(bt_path: Path, processed: Path) -> Path | None:
    if not bt_path.exists():
        return None
    bt = pd.read_parquet(bt_path)
    score_cols = [col for col in bt.columns if col.startswith("s_")]
    if not score_cols:
        return None
    factor = score_cols[0].replace("s_", "")
    ret_col = f"ret_{factor}"
    if ret_col not in bt.columns:
        return None
    signal = bt[score_cols[0]].dropna()
    fwd = bt[ret_col].shift(-1).reindex(signal.index)
    df = pd.DataFrame({"signal": signal, "fwd": fwd}).dropna()
    if df["signal"].nunique() < 5:
        return None
    try:
        buckets = pd.qcut(df["signal"], 5, labels=["Q1", "Q2", "Q3", "Q4", "Q5"])
    except ValueError:
        return None
    quintile = df.groupby(buckets)["fwd"].mean()
    fig = px.bar(
        quintile,
        x=quintile.index,
        y=quintile.values,
        title=f"{factor} Quintile Payoff",
        color=quintile.values,
        color_continuous_scale="Turbo",
    )
    fig.update_layout(xaxis_title="Quintile", yaxis_title="Forward Return")
    return _save_plot(fig, processed / "quintile_summary.png")


def _baseline_macro(bt: pd.DataFrame) -> dict[str, float]:
    if bt.empty:
        return get_macro_features(pd.Timestamp.today())
    last_date = pd.Timestamp(bt.index[-1])
    return get_macro_features(last_date)


def _latest_nlp(processed: Path) -> dict[str, float]:
    llm_path = processed / "llm_policy_log.csv"
    if not llm_path.exists():
        return {"nlp_regime": 0.0, "inflation_mention": 0.0, "growth_mention": 0.0}
    df = pd.read_csv(llm_path)
    if df.empty:
        return {"nlp_regime": 0.0, "inflation_mention": 0.0, "growth_mention": 0.0}
    latest = df.iloc[-1]
    return {
        "nlp_regime": float(latest.get("nlp_regime", 0.0)),
        "inflation_mention": float(latest.get("inflation_mention", 0.0)),
        "growth_mention": float(latest.get("growth_mention", 0.0)),
    }


def _sensitivity_chart(processed: Path) -> Path | None:
    bt_path = processed / "backtest.parquet"
    if not bt_path.exists():
        return None
    bt = pd.read_parquet(bt_path)
    macro = _baseline_macro(bt)
    nlp = _latest_nlp(processed)
    base_weights, _ = _deterministic_policy(macro, nlp)
    base_series = pd.Series(base_weights, dtype=float)
    sensitivities = {}
    for macro_label, delta in SENS_PRESET.items():
        macro_key = MACRO_KEY_MAP[macro_label]
        scenario = macro.copy()
        scenario[macro_key] = scenario.get(macro_key, 0.0) + delta
        mri_adjust = (
            (0.1 * delta if macro_label == "CPI" else 0.0)
            + (-0.2 * delta if macro_label == "UNRATE" else 0.0)
            + (0.35 * delta if macro_label == "T10Y2Y" else 0.0)
        )
        scenario["mri"] = scenario.get("mri", 0.0) + mri_adjust
        new_weights, _ = _deterministic_policy(scenario, nlp)
        new_series = pd.Series(new_weights, dtype=float).reindex(base_series.index).fillna(0.0)
        sensitivities[macro_label] = float((new_series - base_series).abs().sum() / 2.0)
    sens_series = pd.Series(sensitivities)
    fig = px.bar(
        sens_series.sort_values(ascending=True),
        x=sens_series.sort_values(ascending=True).values,
        y=sens_series.sort_values(ascending=True).index,
        orientation="h",
        title="Simulator Sensitivity",
        color=sens_series.sort_values(ascending=True).values,
        color_continuous_scale="Plasma",
    )
    fig.update_layout(xaxis_title="Turnover Impact", yaxis_title="")
    return _save_plot(fig, processed / "simulator_sensitivity.png")


def _add_image_slide(prs: Presentation, title: str, image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.5)).text = title
    slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.0), width=Inches(9))
    return slide


def _add_static_screenshot(prs: Presentation, title: str, path: Path, footer: str) -> None:
    if not path.exists():
        return
    slide = _add_image_slide(prs, title, path)
    _add_footer(slide, footer)


def _add_attribution_slide(prs: Presentation, processed: Path, footer_text: str) -> None:
    ic_path = processed / "ic_summary.csv"
    bt_path = processed / "backtest.parquet"
    if not ic_path.exists() or not bt_path.exists():
        return
    ic_df = pd.read_csv(ic_path).head(5)
    if ic_df.empty:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Attribution Summary"
    rows, cols = ic_df.shape
    table_shape = slide.shapes.add_table(
        rows + 1, cols, Inches(0.5), Inches(1.0), Inches(5), Inches(3)
    )
    table = table_shape.table
    for j, column in enumerate(ic_df.columns):
        table.cell(0, j).text = column
    for i, (_, row) in enumerate(ic_df.iterrows(), start=1):
        for j, value in enumerate(row):
            table.cell(i, j).text = f"{value:.3f}" if isinstance(value, float) else str(value)
    quintile_path = _quintile_chart(bt_path, processed)
    if quintile_path:
        slide.shapes.add_picture(str(quintile_path), Inches(5.7), Inches(1.0), width=Inches(4))
    _add_footer(slide, footer_text)


def main() -> None:
    processed = loader.processed_path()
    cfg = loader._load_config(loader.DEFAULT_CONFIG)
    footer_text = (
        f"Generated with MacroTone v3 | Costs {cfg.portfolio.cost_bps:.0f} bps | "
        f"Benchmark {cfg.project.benchmark}"
    )
    prs = Presentation()
    images: list[tuple[str, Path | None]] = [
        ("Tearsheet Snapshot", processed / "tearsheet.png"),
        ("Equity Curve", processed / "equity_curve.png"),
        ("Rolling Sharpe", processed / "rolling_sharpe.png"),
        ("IC Summary", _ic_chart(processed)),
        ("Factor Attribution", _attribution_chart(processed)),
        ("Simulator Sensitivity Tornado", _sensitivity_chart(processed)),
    ]
    _add_attribution_slide(prs, processed, footer_text)
    for title, image in images:
        if image and Path(image).exists():
            slide = _add_image_slide(prs, title, Path(image))
            _add_footer(slide, footer_text)
    screenshots = {
        "UI Overview": Path("docs/ui/v3_overview.png"),
        "Simulator Presets": Path("docs/ui/v3_simulator.png"),
        "Diagnostics Snapshot": Path("docs/ui/v3_diagnostics.png"),
    }
    for title, path in screenshots.items():
        _add_static_screenshot(prs, title, path, footer_text)
    output = processed / "MacroTone_Tearsheet_v3.pptx"
    prs.save(output)
    print(f"Saved presentation -> {output}")


if __name__ == "__main__":
    main()
